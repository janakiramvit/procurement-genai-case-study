"""
Build the RAG vector store from data/KnowledgeBase/.

Parses PDF / DOCX / PPTX files, chunks the text, embeds each chunk with
OpenAI's text-embedding-3-small, and writes:
  - api/_store/chunks.json     (chunk text + source metadata)
  - api/_store/embeddings.npy  (float32 array, L2-normalized, one row per chunk)

Design choice: chunking is a simple character-based sliding window (~600
tokens/2400 chars, ~100 token/400 char overlap), snapping to whitespace to
avoid mid-word cuts. This avoids a tiktoken dependency for a corpus this
small (~16 documents). At query time we do brute-force numpy cosine
similarity (see api/agents/rag_agent.py) instead of a vector DB -- fine for
a few hundred chunks, documented as a scale-limited tradeoff in DESIGN.md.

The UNGM UNSPSC xlsx is intentionally NOT ingested here: it's a structured
code/title lookup table (13k rows), not prose, so it's loaded into DuckDB
as a queryable table in build_duckdb.py instead -- exact/LIKE lookups are
far more reliable than embedding similarity for taxonomy code lookups.
"""

import json
import os
from io import BytesIO
from pathlib import Path

import numpy as np
from openai import OpenAI
from pptx import Presentation
from pypdf import PdfReader
import docx

ROOT = Path(__file__).resolve().parent.parent
KB_DIR = ROOT / "data" / "KnowledgeBase"
STORE_DIR = ROOT / "api" / "_store"
STORE_DIR.mkdir(parents=True, exist_ok=True)

EMBEDDING_MODEL = "text-embedding-3-small"
CHUNK_SIZE = 2400
CHUNK_OVERLAP = 400
EMBED_BATCH = 100

client = OpenAI()


def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP):
    text = text.strip()
    if not text:
        return []
    chunks = []
    start = 0
    n = len(text)
    while start < n:
        end = min(start + chunk_size, n)
        if end < n:
            # snap to the nearest whitespace so we don't cut mid-word
            snap = text.rfind(" ", start, end)
            if snap > start:
                end = snap
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        if end >= n:
            break
        start = max(end - overlap, start + 1)
    return chunks


def _repair_offbyone_catalog_ref(data: bytes) -> bytes:
    """One KnowledgeBase file (Pharma_UNSPSC_Thresholds.pdf) ships with a corrupted
    trailer/catalog: /Root points to an object number one past the last object
    actually defined, and the Catalog's own /Pages self-references instead of
    pointing at the real Pages tree. Both are off-by-one errors from whatever
    generated the file. Patch them generically: if /Root N 0 R points past the
    highest defined object, and object N-1 is a self-referencing Catalog, fix both."""
    import re

    root_match = re.search(rb"/Root (\d+) 0 R", data)
    if not root_match:
        return data
    root_num = int(root_match.group(1))
    max_obj = max((int(m) for m in re.findall(rb"(\d+) 0 obj", data)), default=0)
    if root_num <= max_obj:
        return data  # root reference is valid, nothing to repair
    fixed_root_num = root_num - 1
    catalog_pattern = f"/Type /Catalog /Pages {fixed_root_num} 0 R".encode()
    if catalog_pattern not in data:
        return data  # doesn't match the known corruption shape, leave as-is
    data = data.replace(f"/Root {root_num} 0 R".encode(), f"/Root {fixed_root_num} 0 R".encode())
    pages_num = fixed_root_num - 1
    data = data.replace(catalog_pattern, f"/Type /Catalog /Pages {pages_num} 0 R".encode())
    return data


def extract_pdf(path: Path):
    """Returns list of (page_number, text)."""
    try:
        reader = PdfReader(str(path))
        _ = len(reader.pages)  # force trailer/catalog resolution now, so we can fall back below
    except Exception:
        repaired = _repair_offbyone_catalog_ref(path.read_bytes())
        reader = PdfReader(BytesIO(repaired))
    pages = []
    for i, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            pages.append((i, text))
    return pages


def extract_docx(path: Path):
    d = docx.Document(str(path))
    paras = [p.text for p in d.paragraphs if p.text.strip()]
    return "\n".join(paras)


def extract_pptx(path: Path):
    """Returns list of (slide_number, text)."""
    prs = Presentation(str(path))
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                texts.append(shape.text.strip())
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells)
                    if row_text.strip(" |"):
                        texts.append(row_text)
        if texts:
            slides.append((i, "\n".join(texts)))
    return slides


def build_records():
    records = []  # {source, location, text}
    files = sorted(KB_DIR.iterdir())
    for path in files:
        suffix = path.suffix.lower()
        try:
            if suffix == ".pdf":
                for page_num, text in extract_pdf(path):
                    for chunk in chunk_text(text):
                        records.append({"source": path.name, "location": f"page {page_num}", "text": chunk})
            elif suffix == ".docx":
                text = extract_docx(path)
                for chunk in chunk_text(text):
                    records.append({"source": path.name, "location": "document", "text": chunk})
            elif suffix == ".pptx":
                for slide_num, text in extract_pptx(path):
                    for chunk in chunk_text(text):
                        records.append({"source": path.name, "location": f"slide {slide_num}", "text": chunk})
            elif suffix == ".xlsx":
                continue  # handled as structured data in build_duckdb.py
            else:
                print(f"skipping unsupported file: {path.name}")
        except Exception as e:
            print(f"WARNING: failed to parse {path.name} ({e.__class__.__name__}: {e}) -- skipping this file")
    return records


def embed_all(records):
    vectors = []
    for i in range(0, len(records), EMBED_BATCH):
        batch = records[i : i + EMBED_BATCH]
        inputs = [r["text"] for r in batch]
        resp = client.embeddings.create(model=EMBEDDING_MODEL, input=inputs)
        for item in resp.data:
            vectors.append(item.embedding)
        print(f"embedded {min(i + EMBED_BATCH, len(records))}/{len(records)} chunks")
    arr = np.array(vectors, dtype=np.float32)
    # L2-normalize so cosine similarity == dot product at query time
    norms = np.linalg.norm(arr, axis=1, keepdims=True)
    norms[norms == 0] = 1.0
    arr = arr / norms
    return arr


def main():
    if not os.environ.get("OPENAI_API_KEY"):
        raise SystemExit("OPENAI_API_KEY not set — create .env.local or export it before running this script")

    print(f"scanning {KB_DIR} ...")
    records = build_records()
    print(f"built {len(records)} chunks from {len(list(KB_DIR.iterdir()))} files")

    embeddings = embed_all(records)

    chunks_path = STORE_DIR / "chunks.json"
    embeddings_path = STORE_DIR / "embeddings.npy"

    with open(chunks_path, "w") as f:
        json.dump(records, f)
    np.save(embeddings_path, embeddings)

    print(f"wrote {chunks_path} ({chunks_path.stat().st_size / 1024:.1f} KB)")
    print(f"wrote {embeddings_path} ({embeddings_path.stat().st_size / 1024:.1f} KB), shape={embeddings.shape}")


if __name__ == "__main__":
    main()
